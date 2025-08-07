# main.py - FINAL CORRECTED VERSION WITH ZIP SUPPORT
import os
import requests
import asyncio
import hashlib
import time
import logging
import shutil
from contextlib import asynccontextmanager
from typing import List, Tuple, Dict, Optional
import pickle
import torch
import uuid
import re
# <<< ADDED: New imports for ZIP file handling >>>
import zipfile
import io

from fastapi import FastAPI, Request, HTTPException
from pydantic import BaseModel
from dotenv import load_dotenv
from openai import AsyncOpenAI
from sentence_transformers import SentenceTransformer
import redis.asyncio as redis

# Import for ColBERT
from ragatouille import RAGPretrainedModel

# --- Document Processing Imports ---
import fitz
import pandas as pd
from docx import Document
from pptx import Presentation
# --- End Document Processing Imports ---

import numpy as np
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Qdrant imports
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

# BM25 for lexical search
from rank_bm25 import BM25Okapi

# --- GPU: Detect device at the very start ---
DEVICE = "cuda" if torch.cuda.is_available() else "cpu"

# --- Configuration (No changes) ---
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
EXPECTED_TOKEN = os.getenv("EXPECTED_TOKEN")
REDIS_URL = os.getenv("REDIS_URL")
QDRANT_URL = os.getenv("QDRANT_URL", "http://localhost:6333")

if not all([OPENAI_API_KEY, EXPECTED_TOKEN, REDIS_URL]):
    raise RuntimeError("Required environment variables are missing.")

EMBEDDING_MODEL_NAME = "BAAI/bge-large-en-v1.5"
EMBEDDING_MODEL_DIMENSION = 1024
RERANKER_MODEL_NAME = "colbert-ir/colbertv2.0"
LLM_MODEL_NAME = "gpt-4o-mini"
TEMP_DIR = "temp_documents"
LOG_DIR = "logs"
os.makedirs(LOG_DIR, exist_ok=True)

# --- Logging Setup (No changes) ---
request_logger = logging.getLogger("api_requests")
request_logger.setLevel(logging.INFO)
request_handler = logging.FileHandler(os.path.join(LOG_DIR, "api_requests.log"), encoding='utf-8')
request_handler.setFormatter(logging.Formatter("%(asctime)s - %(message)s"))
request_logger.addHandler(request_handler)

app_logger = logging.getLogger("rag_system")
app_logger.setLevel(logging.INFO)
app_handler = logging.FileHandler(os.path.join(LOG_DIR, "rag_system.log"), encoding='utf-8')
app_handler.setFormatter(logging.Formatter("%(asctime)s - %(levelname)s - %(message)s"))
stream_handler = logging.StreamHandler()
app_logger.addHandler(app_handler)
app_logger.addHandler(stream_handler)

app_logger.info(f"--- Using device for models: {DEVICE.upper()} ---")
app_logger.info(f"--- Embedding Model: {EMBEDDING_MODEL_NAME} ---")
app_logger.info(f"--- Reranker Model: {RERANKER_MODEL_NAME} ---")

# --- Simple Text Preprocessor (No changes) ---
class SimpleTextProcessor:
    def __init__(self):
        self.stop_words = {'i', 'me', 'my', 'myself', 'we', 'our', 'ours', 'ourselves', 'you', 'your', 'yours', 'yourself', 'yourselves', 'he', 'him', 'his', 'himself', 'she', 'her', 'hers', 'herself', 'it', 'its', 'itself', 'they', 'them', 'their', 'theirs', 'themselves', 'what', 'which', 'who', 'whom', 'this', 'that', 'these', 'those', 'am', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'having', 'do', 'does', 'did', 'doing', 'a', 'an', 'the', 'and', 'but', 'if', 'or', 'because', 'as', 'until', 'while', 'of', 'at', 'by', 'for', 'with', 'through', 'during', 'before', 'after', 'above', 'below', 'up', 'down', 'in', 'out', 'on', 'off', 'over', 'under', 'again', 'further', 'then', 'once', 'here', 'there', 'when', 'where', 'why', 'how', 'all', 'any', 'both', 'each', 'few', 'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not', 'only', 'own', 'same', 'so', 'than', 'too', 'very', 's', 't', 'can', 'will', 'just', 'don', 'should', 'now', 'would', 'could', 'may', 'might', 'must', 'shall'}
    
    def tokenize(self, text: str) -> List[str]:
        text = text.lower()
        text = re.sub(r'[^\w\s]', ' ', text)
        words = text.split()
        return [word for word in words if len(word) > 2 and word not in self.stop_words]

text_processor = SimpleTextProcessor()

# --- Enhanced RAG System Class (No changes) ---
class EnhancedHybridRAGSystem:
    def __init__(self, embedding_model: SentenceTransformer, reranker: RAGPretrainedModel,
                 openai_client: AsyncOpenAI, qdrant_client: QdrantClient):
        self.embedding_model = embedding_model
        self.reranker = reranker
        self.openai_client = openai_client
        self.qdrant_client = qdrant_client
        self.documents: List[str] = []
        self.collection_name = None
        self.bm25_index = None
        
    def is_ready(self) -> bool:
        return (self.collection_name is not None and self.bm25_index is not None and len(self.documents) > 0)
    
    def preprocess_text(self, text: str) -> List[str]:
        return text_processor.tokenize(text)
    
    async def create_qdrant_collection(self, collection_name: str, vector_size: int = EMBEDDING_MODEL_DIMENSION):
        try:
            collections_response = self.qdrant_client.get_collections()
            existing_names = [col.name for col in collections_response.collections]
            
            if collection_name in existing_names:
                app_logger.info(f"Collection '{collection_name}' already exists. Verifying vector size.")
                self.collection_name = collection_name
                
                response_object = self.qdrant_client.get_collection(collection_name=collection_name)
                collection_info = response_object.result if hasattr(response_object, 'result') else response_object
                
                vectors_config_params = collection_info.config.params.vectors
                
                if isinstance(vectors_config_params, VectorParams):
                    existing_size = vectors_config_params.size
                else: 
                    existing_size = vectors_config_params[''].size

                if existing_size != vector_size:
                    app_logger.error(f"CRITICAL MISMATCH: Existing collection '{collection_name}' vector size is {existing_size}, config requires {vector_size}.")
                    raise RuntimeError(f"Qdrant collection '{collection_name}' has incorrect vector dimensions.")
                return
            
            app_logger.info(f"Creating new Qdrant collection: '{collection_name}' with vector size {vector_size}")
            self.qdrant_client.create_collection(
                collection_name=collection_name,
                vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
            )
            self.collection_name = collection_name
        except Exception as e:
            app_logger.error(f"Error creating or verifying Qdrant collection: {e}")
            raise

    async def setup_indexes(self, documents: List[str], doc_hash: str):
        app_logger.info(f"Setting up indexes for {len(documents)} document chunks...")
        self.documents = documents
        collection_name = f"doc_{doc_hash[:12]}"
        await self.create_qdrant_collection(collection_name)
        app_logger.info("Creating embeddings...")
        embeddings = self.embedding_model.encode(documents, convert_to_tensor=True, device=DEVICE, show_progress_bar=True, batch_size=32).cpu().numpy()
        points = [PointStruct(id=str(uuid.uuid4()), vector=embedding.tolist(), payload={"text": doc, "chunk_id": i, "doc_hash": doc_hash}) for i, (doc, embedding) in enumerate(zip(documents, embeddings))]
        batch_size = 100
        for i in range(0, len(points), batch_size):
            self.qdrant_client.upsert(collection_name=self.collection_name, points=points[i:i + batch_size], wait=True)
        app_logger.info(f"Indexed {len(points)} vectors in Qdrant collection: {self.collection_name}")
        app_logger.info("Building BM25 index...")
        try:
            tokenized_docs = [self.preprocess_text(doc) for doc in documents]
            self.bm25_index = BM25Okapi(tokenized_docs)
            app_logger.info("BM25 indexing complete.")
        except Exception as e:
            app_logger.error(f"BM25 indexing failed: {e}")
            self.bm25_index = None
        self.collection_name = collection_name
    
    def hybrid_retrieval(self, question: str, top_k: int = 15) -> List[Tuple[str, int]]:
        if not self.is_ready():
            raise RuntimeError("Indexes not ready for retrieval.")
        query_embedding = self.embedding_model.encode([question], convert_to_tensor=True, device=DEVICE).cpu().numpy()[0]
        vector_results = self.qdrant_client.search(collection_name=self.collection_name, query_vector=query_embedding.tolist(), limit=top_k, with_payload=True)
        bm25_results = []
        if self.bm25_index:
            try:
                tokenized_query = self.preprocess_text(question)
                bm25_scores = self.bm25_index.get_scores(tokenized_query)
                top_indices = np.argsort(bm25_scores)[::-1][:top_k]
                bm25_results = [{'chunk_id': int(idx), 'score': float(bm25_scores[idx])} for idx in top_indices if bm25_scores[idx] > 0]
            except Exception as e:
                app_logger.warning(f"BM25 search failed: {e}")
        combined_results = {}
        k_rrf = 60
        for rank, result in enumerate(vector_results):
            chunk_id = result.payload['chunk_id']
            rrf_score = 1.0 / (k_rrf + rank + 1)
            combined_results[chunk_id] = {'text': result.payload['text'], 'rrf_score': rrf_score}
        for rank, result in enumerate(bm25_results):
            chunk_id = result['chunk_id']
            rrf_score = 1.0 / (k_rrf + rank + 1)
            if chunk_id in combined_results:
                combined_results[chunk_id]['rrf_score'] += rrf_score
            else:
                if 0 <= chunk_id < len(self.documents):
                    combined_results[chunk_id] = {'text': self.documents[chunk_id], 'rrf_score': rrf_score}
        sorted_results = sorted(combined_results.values(), key=lambda x: x['rrf_score'], reverse=True)
        return [(item['text'], i) for i, item in enumerate(sorted_results[:top_k])]

    def rerank_with_colbert(self, question: str, retrieved_docs: List[Tuple[str, int]], top_k: int) -> List[str]:
        if not retrieved_docs:
            return []
        doc_texts = [doc for doc, idx in retrieved_docs]
        app_logger.info(f"Re-ranking {len(doc_texts)} documents with ColBERT...")
        reranked_docs = self.reranker.rerank(query=question, documents=doc_texts, k=top_k)
        return [doc['content'] for doc in reranked_docs]
    
    async def _get_single_answer(self, question: str) -> str:
        try:
            retrieved_docs_with_metadata = self.hybrid_retrieval(question, top_k=15)
            reranked_docs = self.rerank_with_colbert(question, retrieved_docs_with_metadata, top_k=7)
            context = "\n\n---\n\n".join(reranked_docs)
            system_prompt = "You are an expert analysis system. Your task is to answer the user's question in a single, sharp, and precise sentence based on the provided context. If the context is not relevant, use your own knowledge to answer. Never say you don't have enough information. Be direct and to the point."
            user_prompt = f"Context:\n{context}\n\nQuestion: {question}\n\nAnswer in one sentence with all necessary details:"
            response = await self.openai_client.chat.completions.create(model=LLM_MODEL_NAME, messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}], temperature=0.1, max_tokens=1000)
            return response.choices[0].message.content
        except Exception as e:
            app_logger.error(f"Error answering question '{question[:50]}...': {e}", exc_info=True)
            return f"Error: Could not process the answer. Details: {str(e)}"

    async def answer_queries_parallel(self, questions: List[str]) -> List[str]:
        if not self.is_ready():
            raise RuntimeError("RAG system not ready")
        app_logger.info(f"Processing {len(questions)} questions in parallel...")
        tasks = [self._get_single_answer(q) for q in questions]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        final_answers = [str(res) if not isinstance(res, Exception) else f"Error: {res}" for res in results]
        app_logger.info("All questions processed.")
        return final_answers

# --- FastAPI Lifespan & State (No changes) ---
app_state: Dict[str, any] = {}

@asynccontextmanager
async def lifespan(app: FastAPI):
    app_logger.info("--- RAG Server Starting Up ---")
    os.makedirs(TEMP_DIR, exist_ok=True)
    app_state["redis"] = redis.from_url(REDIS_URL, encoding="utf-8", decode_responses=False)
    app_state["qdrant_client"] = QdrantClient(url=QDRANT_URL)
    app_logger.info(f"Loading embedding model: {EMBEDDING_MODEL_NAME}...")
    app_state["embedding_model"] = SentenceTransformer(EMBEDDING_MODEL_NAME, device=DEVICE)
    app_logger.info(f"Loading ColBERT reranker model: {RERANKER_MODEL_NAME}...")
    app_state["reranker"] = RAGPretrainedModel.from_pretrained(RERANKER_MODEL_NAME)
    app_state["openai_client"] = AsyncOpenAI(api_key=OPENAI_API_KEY)
    app_logger.info("--- RAG System Ready ---")
    yield
    app_logger.info("--- Server Shutting Down ---")
    await app_state["redis"].close()
    app_state.clear()

app = FastAPI(title="Enhanced RAG System with ColBERT Reranker", lifespan=lifespan)

# --- API Models (No changes) ---
class HackRxRequest(BaseModel):
    documents: str
    questions: List[str]

class HackRxResponse(BaseModel):
    answers: List[str]
    avg_runtime_per_question: float

# --- Document Processing & Caching ---
# <<< MODIFIED: Function signature now takes a url_hint for extension parsing >>>
def extract_text_from_document(file_path: str, url_hint: str) -> str:
    text = ""
    # Use the hint to determine the file type
    file_extension = os.path.splitext(url_hint.split('?')[0])[-1].lower()
    app_logger.info(f"Extracting text from '{os.path.basename(file_path)}' with extension: {file_extension}")

    try:
        if file_extension == '.pdf':
            with fitz.open(file_path) as doc:
                text = "".join(page.get_text() for page in doc)
        elif file_extension == '.docx':
            doc = Document(file_path)
            text = "\n".join(para.text for para in doc.paragraphs)
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif file_extension == '.xlsx':
            xls = pd.ExcelFile(file_path)
            # Combine all sheets into one text block
            all_sheets_text = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                all_sheets_text.append(f"--- Sheet: {sheet_name} ---\n{df.to_string()}")
            text = "\n\n".join(all_sheets_text)
        else:
            # Silently ignore unsupported files in a zip, or log a warning
            app_logger.warning(f"Unsupported file type '{file_extension}' found in ZIP. Skipping.")
            return "" # Return empty string for unsupported types
    except Exception as e:
        app_logger.error(f"Failed to process document {os.path.basename(file_path)}: {e}")
        raise # Re-raise to be handled in the main processing block
    
    return text

# <<< REWRITTEN: This function now handles both single files and ZIP archives >>>
async def get_or_create_enhanced_rag_system(doc_url: str) -> EnhancedHybridRAGSystem:
    doc_identifier = f"{doc_url}-{EMBEDDING_MODEL_NAME}"
    doc_hash = hashlib.sha256(doc_identifier.encode()).hexdigest()
    redis_client: redis.Redis = app_state["redis"]
    
    rag_system = EnhancedHybridRAGSystem(
        embedding_model=app_state["embedding_model"],
        reranker=app_state["reranker"],
        openai_client=app_state["openai_client"],
        qdrant_client=app_state["qdrant_client"]
    )
    
    cache_key = f"rag_system:{doc_hash}:metadata"
    cached_metadata = await redis_client.get(cache_key)
    if cached_metadata:
        app_logger.info(f"Cache HIT for document hash: {doc_hash}")
        metadata = pickle.loads(cached_metadata)
        rag_system.documents = metadata['documents']
        rag_system.collection_name = metadata['collection_name']
        await rag_system.create_qdrant_collection(rag_system.collection_name)
        if metadata.get('bm25_tokenized_docs'):
            try:
                rag_system.bm25_index = BM25Okapi(metadata['bm25_tokenized_docs'])
            except Exception as e:
                app_logger.warning(f"Failed to restore BM25 index: {e}")
        return rag_system
    
    app_logger.info(f"Cache MISS for document hash: {doc_hash}")
    lock_key = f"lock:rag_system:{doc_hash}"
    
    # Define unique paths for temp files and directories to avoid collisions
    temp_single_file_path = os.path.join(TEMP_DIR, f"{doc_hash}.tmp")
    temp_unzip_dir = os.path.join(TEMP_DIR, f"{doc_hash}_unzipped")

    async with redis_client.lock(lock_key, timeout=600): # Increased timeout for zip processing
        # Double-check cache after acquiring lock
        if await redis_client.get(cache_key):
             # This block can be simplified as the logic is the same as the initial cache check
            return await get_or_create_enhanced_rag_system(doc_url)

        app_logger.info(f"Processing document from URL: {doc_url}")
        try:
            doc_response = requests.get(doc_url, timeout=120)
            doc_response.raise_for_status()
            
            start_time = time.time()
            file_extension = os.path.splitext(doc_url.split('?')[0])[-1].lower()
            combined_text = ""

            if file_extension == '.zip':
                app_logger.info("Detected ZIP archive. Extracting and processing all files...")
                os.makedirs(temp_unzip_dir, exist_ok=True)
                
                with zipfile.ZipFile(io.BytesIO(doc_response.content)) as z:
                    z.extractall(temp_unzip_dir)
                
                all_texts = []
                for root, _, files in os.walk(temp_unzip_dir):
                    for filename in files:
                        file_path = os.path.join(root, filename)
                        try:
                            # Pass filename as the hint to parse the correct extension
                            text_from_file = extract_text_from_document(file_path, filename)
                            if text_from_file.strip():
                                all_texts.append(text_from_file)
                        except Exception as e:
                            app_logger.warning(f"Could not process file '{filename}' in zip: {e}")
                
                # Combine text from all files with a clear separator
                combined_text = "\n\n--- NEW DOCUMENT ---\n\n".join(all_texts)
            else:
                # Logic for single, non-zip files
                with open(temp_single_file_path, "wb") as f:
                    f.write(doc_response.content)
                combined_text = extract_text_from_document(temp_single_file_path, doc_url)

            if not combined_text.strip():
                app_logger.warning("Document is empty or no supported files found after extraction.")
                return rag_system
            
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=512, chunk_overlap=100, separators=["\n\n", "\n", ". ", "!", "?", ",", " ", ""])
            documents = text_splitter.split_text(combined_text)
            
            if not documents:
                app_logger.warning("No document chunks created after splitting.")
                return rag_system
            
            app_logger.info(f"Created {len(documents)} document chunks from all processed files.")
            await rag_system.setup_indexes(documents, doc_hash)
            
            processing_time = time.time() - start_time
            app_logger.info(f"Document processing completed in {processing_time:.2f} seconds.")
            
            metadata = {
                'documents': documents, 
                'collection_name': rag_system.collection_name, 
                'bm25_tokenized_docs': [rag_system.preprocess_text(doc) for doc in documents] if rag_system.bm25_index else None,
                'doc_hash': doc_hash
            }
            await redis_client.setex(cache_key, 86400, pickle.dumps(metadata))
            app_logger.info(f"Document {doc_hash} processed and cached successfully.")

        finally:
            # Securely clean up temp files and directories
            if os.path.exists(temp_single_file_path):
                os.remove(temp_single_file_path)
            if os.path.exists(temp_unzip_dir):
                shutil.rmtree(temp_unzip_dir)
    
    return rag_system

# --- API Endpoint & Health Check (No changes) ---
@app.post("/hackrx/run", response_model=HackRxResponse)
async def run_submission(payload: HackRxRequest, request: Request):
    start_time = time.time()
    request_logger.info(f"POST /hackrx/run BODY: {payload.model_dump_json()}")
    auth_header = request.headers.get("Authorization")
    if not auth_header or auth_header != f"Bearer {EXPECTED_TOKEN}":
        raise HTTPException(status_code=401, detail="Invalid auth credentials")
    try:
        rag_system = await get_or_create_enhanced_rag_system(doc_url=payload.documents)
        if not rag_system.is_ready():
            raise HTTPException(status_code=500, detail="RAG system could not be initialized for the document.")
        answers = await rag_system.answer_queries_parallel(questions=payload.questions)
        process_time = time.time() - start_time
        avg_runtime = process_time / len(payload.questions) if payload.questions else 0
        app_logger.info(f"Request completed. Total time: {process_time:.2f}s, Avg time/q: {avg_runtime:.2f}s")
        return HackRxResponse(answers=answers, avg_runtime_per_question=avg_runtime)
    except Exception as e:
        app_logger.error(f"An unexpected error occurred during request processing: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/health")
async def health_check():
    return {
        "status": "healthy",
        "device": DEVICE,
        "models_loaded": "embedding_model" in app_state and "reranker" in app_state,
        "embedding_model": EMBEDDING_MODEL_NAME,
        "reranker_model": RERANKER_MODEL_NAME,
        "qdrant_url": QDRANT_URL
    }